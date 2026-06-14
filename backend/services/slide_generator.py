"""
Slide Deck Generator Service
=============================
Generates professional slide decks from lesson plan content.

Three-phase pipeline:
  1. Content generation (Gemini 2.5 Flash text) — structured JSON with layouts
  2. Image generation (Gemini 2.5 Flash Image via google-genai) — themed graphics
  3. PPTX assembly (python-pptx) — builds the PowerPoint file

Usage:
    from backend.services.slide_generator import generate_slide_content, generate_slide_images, assemble_pptx

    content = generate_slide_content(text, subject, grade, title, api_key)
    images = generate_slide_images(content["slides"], content["theme"], api_key)
    assemble_pptx(content["slides"], content["theme"], content["title"], images, filepath)
"""

import json
import logging
import os
from io import BytesIO

import sentry_sdk

logger = logging.getLogger(__name__)

# ── Default color fallback (used only if AI doesn't provide colors) ────

DEFAULT_THEME = {
    "primary_color": "#1a56db",
    "secondary_color": "#60a5fa",
    "accent": "#dbeafe",
}


# ── Phase 1: Content generation ─────────────────────────────────────────

def generate_slide_content(content, subject, grade, title, api_key,
                           lesson_plan=None, global_ai_notes="", instructions="",
                           slide_count=10, deck_format="detailed", template="academic"):
    """Generate structured slide content from source material.

    The AI decides colors, layouts, and visual style based on the content.
    Teachers optionally provide style instructions.

    Args:
        deck_format: "detailed" (full text, standalone) or "presenter" (key talking points only)

    Returns dict with: title, theme, slides[]
    """
    if not content and not lesson_plan:
        raise ValueError("Provide content or lesson_plan to generate slides.")

    from backend.retry import with_retry

    format_instruction = ""
    if deck_format == "presenter":
        format_instruction = "Create PRESENTER SLIDES: clean, visual slides with only key talking points (2-3 words per bullet). The speaker notes should contain the full detail. Slides should be visually clean with minimal text."
    else:
        format_instruction = "Create a DETAILED DECK: comprehensive slides with full text and details, suitable for emailing or reading standalone without a presenter."

    prompt_parts = []
    prompt_parts.append("You are an expert " + subject + " teacher creating a slide deck for grade " + grade + " students.")
    prompt_parts.append("Create exactly " + str(slide_count) + " slides.")
    prompt_parts.append("")
    prompt_parts.append(format_instruction)

    if global_ai_notes:
        prompt_parts.append("")
        prompt_parts.append("=== TEACHER INSTRUCTIONS (MUST FOLLOW) ===")
        prompt_parts.append(global_ai_notes)
        prompt_parts.append("=== END TEACHER INSTRUCTIONS ===")

    prompt_parts.append("")
    prompt_parts.append("Return ONLY valid JSON with this structure:")
    prompt_parts.append('{')
    prompt_parts.append('  "title": "Deck Title",')
    prompt_parts.append('  "theme": {')
    prompt_parts.append('    "primary_color": "#hex — a single accent color appropriate to the SUBJECT. The deck design (fonts, layout, spacing) is fixed by a professional template; you ONLY choose this one accent color.",')
    prompt_parts.append('    "secondary_color": "#hex — optional complementary accent"')
    prompt_parts.append('  },')
    prompt_parts.append('  "slides": [')
    prompt_parts.append('    {')
    prompt_parts.append('      "layout": "title|content|image_focus|two_column|key_concept|section_divider",')
    prompt_parts.append('      "title": "Slide Title",')
    prompt_parts.append('      "subtitle": "Only for title layout",')
    prompt_parts.append('      "bullets": ["Point 1", "Point 2"],')
    prompt_parts.append('      "content": "For key_concept layout - one big idea",')
    prompt_parts.append('      "left_title": "For two_column",')
    prompt_parts.append('      "left_bullets": ["Left 1"],')
    prompt_parts.append('      "right_title": "For two_column",')
    prompt_parts.append('      "right_bullets": ["Right 1"],')
    prompt_parts.append('      "caption": "For image_focus layout",')
    prompt_parts.append('      "image_prompt": "Describe an illustration for this slide. Be specific about the subject matter. Do NOT include text in the image.",')
    prompt_parts.append('      "speaker_notes": "What the teacher should say during this slide"')
    prompt_parts.append('    }')
    prompt_parts.append('  ]')
    prompt_parts.append('}')
    prompt_parts.append("")
    prompt_parts.append("Layout rules:")
    prompt_parts.append("- Slide 1 MUST be layout 'title' with the deck title and subtitle")
    prompt_parts.append("- Use 'content' for most slides (title + 3-5 bullet points + image)")
    prompt_parts.append("- Use 'key_concept' for 1-2 important ideas (large centered text)")
    prompt_parts.append("- Use 'two_column' for comparisons (Federalists vs Anti-Federalists)")
    prompt_parts.append("- Use 'section_divider' between major topics")
    prompt_parts.append("- Use 'image_focus' when a visual needs to be the main focus")
    prompt_parts.append("- Max 5 bullet points per slide. Keep bullets concise.")
    prompt_parts.append("- Every content and title slide MUST have an image_prompt")
    prompt_parts.append("- image_prompt should describe a flat, educational illustration. NO text in images.")
    prompt_parts.append("- speaker_notes on EVERY slide (2-3 sentences of what the teacher says)")
    prompt_parts.append("")
    prompt_parts.append("Return ONLY valid JSON. No markdown, no code fences.")

    if lesson_plan:
        prompt_parts.append("")
        prompt_parts.append("=== LESSON PLAN ===")
        prompt_parts.append("Title: " + (lesson_plan.get('title', '') or ''))
        if lesson_plan.get('overview'):
            prompt_parts.append("Overview: " + lesson_plan['overview'])
        if lesson_plan.get('objectives'):
            prompt_parts.append("Objectives:")
            for obj in lesson_plan['objectives']:
                prompt_parts.append("  - " + str(obj))
        if lesson_plan.get('vocabulary'):
            prompt_parts.append("Vocabulary: " + ", ".join(lesson_plan['vocabulary']))
        if lesson_plan.get('essential_questions'):
            prompt_parts.append("Essential Questions:")
            for eq in lesson_plan['essential_questions']:
                prompt_parts.append("  - " + str(eq))
        if lesson_plan.get('days'):
            for day in lesson_plan['days']:
                prompt_parts.append("Day " + str(day.get('day', '?')) + ": " + str(day.get('topic', '')))
        prompt_parts.append("=== END LESSON PLAN ===")

    if content:
        prompt_parts.append("")
        prompt_parts.append("=== SOURCE CONTENT ===")
        prompt_parts.append(content[:8000])
        prompt_parts.append("=== END SOURCE CONTENT ===")

    if instructions:
        prompt_parts.append("")
        prompt_parts.append("Additional instructions: " + instructions)

    prompt = "\n".join(prompt_parts)

    from backend.services.llm_adapter import GeminiAdapter, LLMRequest, Message, TextPart
    adapter = GeminiAdapter(api_key=api_key)
    llm_resp = adapter.chat(LLMRequest(
        model="gemini-2.5-flash",
        messages=[Message(role="user", content=[TextPart(text=prompt)])],
        metadata={"feature_label": "generate_slide_content"},
    ))

    response_text = (llm_resp.content_parts[0].text if llm_resp.content_parts else "").strip()
    if response_text.startswith("```"):
        response_text = response_text.split("\n", 1)[1] if "\n" in response_text else response_text[3:]
    if response_text.endswith("```"):
        response_text = response_text[:-3].strip()
    if response_text.startswith("json"):
        response_text = response_text[4:].strip()

    result = json.loads(response_text)

    # Ensure theme exists with required fields (AI should have generated it)
    theme = result.get("theme", {})
    if not theme.get("primary_color"):
        theme["primary_color"] = "#1a56db"
        theme["secondary_color"] = "#60a5fa"
        theme["accent"] = "#dbeafe"
    result["theme"] = theme

    # Build the image style prompt from the AI's chosen theme
    result["theme"]["style_prompt"] = (
        "flat vector illustration, clean minimal educational style, "
        "soft color palette using " + theme.get("primary_color", DEFAULT_THEME["primary_color"]) + " and " + theme.get("secondary_color", DEFAULT_THEME["secondary_color"]) + ", "
        "no text in the image, professional, suitable for a classroom presentation"
    )

    result["template"] = template
    return result


# ── Phase 2: Image generation ───────────────────────────────────────────

def generate_slide_images(slides, theme, api_key, max_images=5):
    """Generate themed images for slides that have image_prompt.

    Uses GeminiAdapter.generate_image for all observability, retry, and
    breaker protection. Reference-image style consistency is preserved:
    the first successful image's bytes + MIME are carried forward as an
    ImagePart passed via reference_images on subsequent calls.

    Args:
        slides: list of slide dicts
        theme: theme dict with style_prompt
        api_key: Gemini API key
        max_images: cap on number of images (cost control)

    Returns:
        dict mapping slide index -> PNG/JPEG image bytes

    Error handling:
        Each image generation is try/excepted individually. If one fails
        (timeout, safety filter, quota, CircuitBreakerError), that slide
        gets no image — the PPTX assembler renders it as text-only. No
        retry on image calls (via ImageRequest(retry=False)).
    """
    import base64
    import pybreaker
    from backend.services.llm_adapter.gemini_adapter import GeminiAdapter
    from backend.services.llm_adapter.types import ImagePart, ImageRequest

    adapter = GeminiAdapter(api_key=api_key)
    style_prompt = theme.get("style_prompt", "flat educational illustration, no text")

    # Find slides that need images
    image_slides = []
    for i, slide in enumerate(slides):
        if slide.get("image_prompt"):
            image_slides.append((i, slide["image_prompt"]))

    image_slides = image_slides[:max_images]

    images = {}
    reference_image_bytes: bytes | None = None
    reference_image_mime: str = "image/png"  # starts as png; SDK response may update

    for idx, (slide_index, image_prompt) in enumerate(image_slides):
        try:
            full_prompt = style_prompt + ". " + image_prompt
            reference_images = None
            if reference_image_bytes is not None:
                b64 = base64.b64encode(reference_image_bytes).decode("ascii")
                reference_images = [ImagePart(
                    url=None,
                    base64=b64,
                    mime_type=reference_image_mime,
                )]
                style_note = (
                    "Generate an illustration in the EXACT same visual style as "
                    "the reference image above. "
                )
                full_prompt = style_note + full_prompt

            response = adapter.generate_image(ImageRequest(
                prompt=full_prompt,
                model="gemini-2.5-flash-image",
                reference_images=reference_images,
                aspect_ratio="16:9",
                metadata={"feature_label": "slide_generator_image"},
            ))

            if response.images:
                image_data = response.images[0]
                images[slide_index] = image_data
                if reference_image_bytes is None:
                    reference_image_bytes = image_data
                    reference_image_mime = response.mime_type
                    logger.info("Style reference image set from slide %d", slide_index)
            # else: empty response — adapter already emitted llm.image.call.blocked.
            # Slide will render text-only in PPTX assembly.
        except Exception as e:  # noqa: BLE001  # broad catch: error is logged
            # Filter CircuitBreakerError out of Sentry to avoid alert fatigue
            # when the breaker stays OPEN across many slides. logger.warning
            # still fires for every failure so the pattern is visible in logs.
            if not isinstance(e, pybreaker.CircuitBreakerError):
                sentry_sdk.capture_exception(e)
            logger.warning(
                "Image generation failed for slide %d (will render text-only): %s",
                slide_index, e,
            )
            # Do NOT retry — $0.04 loss is acceptable vs double-charging

    logger.info("Generated %d/%d slide images", len(images), len(image_slides))
    return images


# ── Phase 3: PPTX assembly ──────────────────────────────────────────────

def _hex_to_rgb(hex_color):
    """Convert hex color string to RGBColor."""
    from pptx.dml.color import RGBColor
    h = hex_color.lstrip('#')
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def _add_title_slide(prs, slide_data, theme, image_bytes=None):
    """Add a title slide with large title, subtitle, and optional background image."""
    from pptx.util import Inches, Pt, Emu
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout

    # Background color bar
    from pptx.util import Inches
    bg_shape = slide.shapes.add_shape(
        1, 0, 0, prs.slide_width, prs.slide_height  # MSO_SHAPE.RECTANGLE = 1
    )
    bg_shape.fill.solid()
    bg_shape.fill.fore_color.rgb = _hex_to_rgb(theme.get("primary_color", "#1a56db"))
    bg_shape.line.fill.background()

    # Add image if available (right side, slightly transparent overlay feel)
    if image_bytes:
        img_stream = BytesIO(image_bytes)
        pic = slide.shapes.add_picture(img_stream, Inches(6.5), Inches(0.5), Inches(6), Inches(6))

    # Title text
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(6), Inches(2.5))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = slide_data.get("title", "")
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = _hex_to_rgb("#ffffff")
    p.alignment = PP_ALIGN.LEFT

    # Subtitle
    if slide_data.get("subtitle"):
        p2 = tf.add_paragraph()
        p2.text = slide_data["subtitle"]
        p2.font.size = Pt(22)
        p2.font.color.rgb = _hex_to_rgb("#e0e0e0")
        p2.alignment = PP_ALIGN.LEFT
        p2.space_before = Pt(16)

    # Speaker notes
    if slide_data.get("speaker_notes"):
        notes_slide = slide.notes_slide
        notes_slide.notes_text_frame.text = slide_data["speaker_notes"]


def _add_content_slide(prs, slide_data, theme, image_bytes=None):
    """Add a content slide with title, bullets, and optional image on right."""
    from pptx.util import Inches, Pt

    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Accent bar at top
    bar = slide.shapes.add_shape(1, 0, 0, prs.slide_width, Inches(0.08))
    bar.fill.solid()
    bar.fill.fore_color.rgb = _hex_to_rgb(theme.get("primary_color", "#1a56db"))
    bar.line.fill.background()

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(7), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = slide_data.get("title", "")
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = _hex_to_rgb(theme.get("primary_color", "#1a56db"))

    # Bullets (left side if image, full width if no image)
    bullet_width = Inches(5.5) if image_bytes else Inches(11)
    bullet_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.6), bullet_width, Inches(5))
    tf = bullet_box.text_frame
    tf.word_wrap = True

    for i, bullet in enumerate(slide_data.get("bullets", [])):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = bullet
        p.font.size = Pt(18)
        p.font.color.rgb = _hex_to_rgb("#374151")
        p.space_before = Pt(12)
        p.level = 0

    # Image on right
    if image_bytes:
        img_stream = BytesIO(image_bytes)
        slide.shapes.add_picture(img_stream, Inches(7), Inches(1.6), Inches(5.5), Inches(5))

    if slide_data.get("speaker_notes"):
        notes_slide = slide.notes_slide
        notes_slide.notes_text_frame.text = slide_data["speaker_notes"]


def _add_key_concept_slide(prs, slide_data, theme, image_bytes=None):
    """Add a key concept slide with large centered text."""
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Light background
    bg = slide.shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    accent = theme.get("accent", "#dbeafe")
    bg.fill.fore_color.rgb = _hex_to_rgb(accent)
    bg.line.fill.background()

    # Centered title
    title_box = slide.shapes.add_textbox(Inches(1.5), Inches(1.5), Inches(10), Inches(1.5))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = slide_data.get("title", "")
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = _hex_to_rgb(theme.get("primary_color", "#1a56db"))
    p.alignment = PP_ALIGN.CENTER

    # Big content text
    if slide_data.get("content"):
        content_box = slide.shapes.add_textbox(Inches(2), Inches(3.5), Inches(9), Inches(2.5))
        tf = content_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = slide_data["content"]
        p.font.size = Pt(24)
        p.font.color.rgb = _hex_to_rgb("#374151")
        p.alignment = PP_ALIGN.CENTER

    if slide_data.get("speaker_notes"):
        notes_slide = slide.notes_slide
        notes_slide.notes_text_frame.text = slide_data["speaker_notes"]


def _add_two_column_slide(prs, slide_data, theme, image_bytes=None):
    """Add a two-column comparison slide."""
    from pptx.util import Inches, Pt

    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Accent bar
    bar = slide.shapes.add_shape(1, 0, 0, prs.slide_width, Inches(0.08))
    bar.fill.solid()
    bar.fill.fore_color.rgb = _hex_to_rgb(theme.get("primary_color", "#1a56db"))
    bar.line.fill.background()

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = slide_data.get("title", "")
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = _hex_to_rgb(theme.get("primary_color", "#1a56db"))

    # Left column
    left_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(5.5), Inches(5))
    tf = left_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = slide_data.get("left_title", "")
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = _hex_to_rgb(theme.get("secondary_color", "#60a5fa"))
    for bullet in slide_data.get("left_bullets", []):
        p2 = tf.add_paragraph()
        p2.text = bullet
        p2.font.size = Pt(16)
        p2.font.color.rgb = _hex_to_rgb("#374151")
        p2.space_before = Pt(8)

    # Divider line
    divider = slide.shapes.add_shape(1, Inches(6.4), Inches(1.6), Inches(0.04), Inches(5))
    divider.fill.solid()
    divider.fill.fore_color.rgb = _hex_to_rgb("#d1d5db")
    divider.line.fill.background()

    # Right column
    right_box = slide.shapes.add_textbox(Inches(6.8), Inches(1.6), Inches(5.5), Inches(5))
    tf = right_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = slide_data.get("right_title", "")
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = _hex_to_rgb(theme.get("secondary_color", "#60a5fa"))
    for bullet in slide_data.get("right_bullets", []):
        p2 = tf.add_paragraph()
        p2.text = bullet
        p2.font.size = Pt(16)
        p2.font.color.rgb = _hex_to_rgb("#374151")
        p2.space_before = Pt(8)

    if slide_data.get("speaker_notes"):
        notes_slide = slide.notes_slide
        notes_slide.notes_text_frame.text = slide_data["speaker_notes"]


def _add_image_focus_slide(prs, slide_data, theme, image_bytes=None):
    """Add a slide with a large centered image and caption."""
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN

    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.3), Inches(11), Inches(0.7))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = slide_data.get("title", "")
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = _hex_to_rgb(theme.get("primary_color", "#1a56db"))
    p.alignment = PP_ALIGN.CENTER

    # Large image
    if image_bytes:
        img_stream = BytesIO(image_bytes)
        slide.shapes.add_picture(img_stream, Inches(1.5), Inches(1.2), Inches(10), Inches(5))

    # Caption
    if slide_data.get("caption"):
        cap_box = slide.shapes.add_textbox(Inches(1.5), Inches(6.5), Inches(10), Inches(0.6))
        tf = cap_box.text_frame
        p = tf.paragraphs[0]
        p.text = slide_data["caption"]
        p.font.size = Pt(14)
        p.font.italic = True
        p.font.color.rgb = _hex_to_rgb("#6b7280")
        p.alignment = PP_ALIGN.CENTER

    if slide_data.get("speaker_notes"):
        notes_slide = slide.notes_slide
        notes_slide.notes_text_frame.text = slide_data["speaker_notes"]


def _add_section_divider_slide(prs, slide_data, theme, image_bytes=None):
    """Add a section divider slide with title and decorative bar."""
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Colored background
    bg = slide.shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = _hex_to_rgb(theme.get("primary_color", "#1a56db"))
    bg.line.fill.background()

    # Decorative bar
    bar = slide.shapes.add_shape(1, Inches(4), Inches(3.2), Inches(5), Inches(0.06))
    bar.fill.solid()
    bar.fill.fore_color.rgb = _hex_to_rgb(theme.get("secondary_color", "#60a5fa"))
    bar.line.fill.background()

    # Section title
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(11), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = slide_data.get("title", "")
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = _hex_to_rgb("#ffffff")
    p.alignment = PP_ALIGN.CENTER

    if slide_data.get("speaker_notes"):
        notes_slide = slide.notes_slide
        notes_slide.notes_text_frame.text = slide_data["speaker_notes"]


# Layout dispatcher
_LAYOUT_HANDLERS = {
    "title": _add_title_slide,
    "content": _add_content_slide,
    "key_concept": _add_key_concept_slide,
    "two_column": _add_two_column_slide,
    "image_focus": _add_image_focus_slide,
    "section_divider": _add_section_divider_slide,
}


def assemble_pptx(slides, theme, title, images, filepath):
    """Assemble a PPTX file from slide data and generated images.

    Args:
        slides: list of slide dicts
        theme: theme dict with colors
        title: deck title
        images: dict mapping slide index → PNG bytes
        filepath: output .pptx path
    """
    from pptx import Presentation

    prs = Presentation()
    prs.slide_width = 12192000   # 13.333 inches in EMU (16:9)
    prs.slide_height = 6858000   # 7.5 inches in EMU

    for i, slide_data in enumerate(slides):
        layout = slide_data.get("layout", "content")
        handler = _LAYOUT_HANDLERS.get(layout, _add_content_slide)
        image_bytes = images.get(i)
        handler(prs, slide_data, theme, image_bytes)

    prs.save(filepath)
    logger.info("Assembled %d-slide PPTX: %s", len(slides), filepath)
