"""Tests for slide deck generation service."""

import json
import os
import tempfile
import pytest
from unittest.mock import patch, MagicMock
from io import BytesIO


SAMPLE_SLIDE_CONTENT = {
    "title": "The Constitution",
    "theme": {
        "primary_color": "#1a56db",
        "secondary_color": "#f59e0b",
        "style_prompt": "flat vector illustration, educational, clean minimal, soft blue and amber palette, no text in images"
    },
    "slides": [
        {
            "layout": "title",
            "title": "The Constitution",
            "subtitle": "Foundations of American Government",
            "image_prompt": "an illustration of the US Constitution document with a quill pen",
            "speaker_notes": "Welcome to our lesson on the Constitution."
        },
        {
            "layout": "content",
            "title": "Three Branches of Government",
            "bullets": [
                "Legislative Branch — makes laws (Congress)",
                "Executive Branch — enforces laws (President)",
                "Judicial Branch — interprets laws (Supreme Court)"
            ],
            "image_prompt": "three pillars representing the branches of government",
            "speaker_notes": "The Constitution divides power into three branches."
        },
        {
            "layout": "key_concept",
            "title": "Separation of Powers",
            "content": "No single branch has too much power. Each branch checks the others.",
            "speaker_notes": "This is the key concept of the lesson."
        },
        {
            "layout": "two_column",
            "title": "Federalists vs Anti-Federalists",
            "left_title": "Federalists",
            "left_bullets": ["Strong central government", "Supported ratification"],
            "right_title": "Anti-Federalists",
            "right_bullets": ["Wanted Bill of Rights", "Feared tyranny"],
            "speaker_notes": "Compare the two sides."
        },
        {
            "layout": "content",
            "title": "The Bill of Rights",
            "bullets": [
                "First 10 amendments to the Constitution",
                "Guarantees individual freedoms",
                "Added to gain Anti-Federalist support"
            ],
            "image_prompt": "a scroll with amendments written on it",
            "speaker_notes": "The Bill of Rights was a compromise."
        },
    ]
}


class TestGenerateSlideContent:
    def test_returns_structured_json(self):
        """Should return slide content with title, theme, and slides array."""
        from backend.services.slide_generator import generate_slide_content
        from backend.services.llm_adapter.types import LLMResponse, TextPart, Usage

        mock_llm_resp = LLMResponse(
            content_parts=[TextPart(text=json.dumps(SAMPLE_SLIDE_CONTENT))],
            tool_calls=[],
            usage=Usage(prompt_tokens=10, completion_tokens=5, cost_usd=0.0),
            finish_reason="stop",
            provider="gemini",
            model="gemini-2.0-flash",
        )

        with patch('backend.services.llm_adapter.gemini_adapter.genai') as mock_genai:
            mock_client = MagicMock()
            mock_client.models.generate_content.return_value = MagicMock(
                text=json.dumps(SAMPLE_SLIDE_CONTENT),
                candidates=[MagicMock(finish_reason=MagicMock(name="STOP"))],
                usage_metadata=MagicMock(prompt_token_count=10, candidates_token_count=5),
            )
            mock_genai.Client.return_value = mock_client

            result = generate_slide_content(
                content="The Constitution establishes three branches...",
                subject="US History",
                grade="8",
                title="The Constitution",
                api_key="fake-key",
            )

        assert "slides" in result
        assert "theme" in result
        assert len(result["slides"]) > 0
        assert result["slides"][0]["layout"] in ("title", "content", "key_concept", "two_column", "image_focus", "section_divider")

    def test_returns_error_without_content(self):
        """Should raise ValueError without content."""
        from backend.services.slide_generator import generate_slide_content

        with pytest.raises(ValueError, match="content"):
            generate_slide_content(content="", subject="Math", grade="7", title="Test", api_key="fake")


class TestGenerateSlideImages:
    def test_generates_images_for_slides_with_prompts(self):
        """Should generate images only for slides that have image_prompt."""
        from backend.services.slide_generator import generate_slide_images

        slides = SAMPLE_SLIDE_CONTENT["slides"]
        theme = SAMPLE_SLIDE_CONTENT["theme"]

        # Create a 1x1 white PNG for the mock
        from PIL import Image
        img = Image.new('RGB', (1024, 576), 'white')
        img_bytes = BytesIO()
        img.save(img_bytes, format='PNG')
        mock_image_data = img_bytes.getvalue()

        mock_part = MagicMock()
        mock_part.inline_data = MagicMock()
        mock_part.inline_data.data = mock_image_data

        mock_response = MagicMock()
        mock_response.candidates = [MagicMock()]
        mock_response.candidates[0].content.parts = [mock_part]

        with patch('backend.services.slide_generator._get_genai_client') as mock_get:
            mock_client = MagicMock()
            mock_client.models.generate_content.return_value = mock_response
            mock_get.return_value = mock_client

            images = generate_slide_images(slides, theme, api_key="fake-key", max_images=5)

        # Only slides with image_prompt get images (3 out of 5 in sample)
        assert len(images) <= 5
        assert len(images) > 0

    def test_respects_max_images_cap(self):
        """Should not generate more images than max_images."""
        from backend.services.slide_generator import generate_slide_images

        slides = [{"layout": "content", "image_prompt": "img " + str(i)} for i in range(20)]
        theme = SAMPLE_SLIDE_CONTENT["theme"]

        from PIL import Image
        img = Image.new('RGB', (100, 56), 'white')
        img_bytes = BytesIO()
        img.save(img_bytes, format='PNG')

        mock_part = MagicMock()
        mock_part.inline_data = MagicMock()
        mock_part.inline_data.data = img_bytes.getvalue()

        mock_response = MagicMock()
        mock_response.candidates = [MagicMock()]
        mock_response.candidates[0].content.parts = [mock_part]

        with patch('backend.services.slide_generator._get_genai_client') as mock_get:
            mock_client = MagicMock()
            mock_client.models.generate_content.return_value = mock_response
            mock_get.return_value = mock_client

            images = generate_slide_images(slides, theme, api_key="fake-key", max_images=5)

        assert len(images) <= 5


class TestAssemblePptx:
    def test_creates_valid_pptx(self):
        """Should create a valid .pptx file from slide data."""
        from backend.services.slide_generator import assemble_pptx

        slides = SAMPLE_SLIDE_CONTENT["slides"]
        theme = SAMPLE_SLIDE_CONTENT["theme"]
        title = SAMPLE_SLIDE_CONTENT["title"]

        filepath = os.path.join(tempfile.mkdtemp(), "test.pptx")
        assemble_pptx(slides, theme, title, {}, filepath)

        assert os.path.exists(filepath)
        assert os.path.getsize(filepath) > 10000  # Non-trivial PPTX

        # Verify it's a valid PPTX by opening it
        from pptx import Presentation
        prs = Presentation(filepath)
        assert len(prs.slides) == len(slides)

    def test_includes_speaker_notes(self):
        """Speaker notes should be present on each slide."""
        from backend.services.slide_generator import assemble_pptx
        from pptx import Presentation

        slides = SAMPLE_SLIDE_CONTENT["slides"]
        theme = SAMPLE_SLIDE_CONTENT["theme"]

        filepath = os.path.join(tempfile.mkdtemp(), "test.pptx")
        assemble_pptx(slides, theme, "Test", {}, filepath)

        prs = Presentation(filepath)
        for i, slide in enumerate(prs.slides):
            notes = slide.notes_slide.notes_text_frame.text if slide.has_notes_slide else ""
            if slides[i].get("speaker_notes"):
                assert slides[i]["speaker_notes"] in notes

    def test_embeds_images(self):
        """Slides with images should have picture shapes."""
        from backend.services.slide_generator import assemble_pptx
        from pptx import Presentation
        from pptx.shapes.picture import Picture
        from PIL import Image

        slides = [SAMPLE_SLIDE_CONTENT["slides"][1]]  # Content slide with image_prompt
        theme = SAMPLE_SLIDE_CONTENT["theme"]

        # Create a test image
        img = Image.new('RGB', (1024, 576), 'blue')
        img_bytes = BytesIO()
        img.save(img_bytes, format='PNG')
        images = {0: img_bytes.getvalue()}

        filepath = os.path.join(tempfile.mkdtemp(), "test.pptx")
        assemble_pptx(slides, theme, "Test", images, filepath)

        prs = Presentation(filepath)
        slide = prs.slides[0]
        has_picture = any(isinstance(shape, Picture) for shape in slide.shapes)
        assert has_picture

    def test_handles_empty_slides(self):
        """Should not crash with empty slides list."""
        from backend.services.slide_generator import assemble_pptx

        filepath = os.path.join(tempfile.mkdtemp(), "test.pptx")
        assemble_pptx([], {"primary_color": "#000"}, "Empty", {}, filepath)

        from pptx import Presentation
        prs = Presentation(filepath)
        assert len(prs.slides) == 0

    def test_all_layout_types(self):
        """Should handle all 6 layout types without crashing."""
        from backend.services.slide_generator import assemble_pptx
        from pptx import Presentation

        slides = [
            {"layout": "title", "title": "Title", "subtitle": "Sub"},
            {"layout": "content", "title": "Content", "bullets": ["A", "B"]},
            {"layout": "image_focus", "title": "Image", "caption": "Caption"},
            {"layout": "two_column", "title": "Compare", "left_title": "L", "left_bullets": ["1"], "right_title": "R", "right_bullets": ["2"]},
            {"layout": "key_concept", "title": "Key", "content": "Big idea"},
            {"layout": "section_divider", "title": "Next Section"},
        ]

        filepath = os.path.join(tempfile.mkdtemp(), "test.pptx")
        assemble_pptx(slides, {"primary_color": "#1a56db", "secondary_color": "#f59e0b"}, "All Layouts", {}, filepath)

        prs = Presentation(filepath)
        assert len(prs.slides) == 6


class TestDeckFormat:
    def test_detailed_format_includes_full_text(self):
        """Detailed format prompt should mention full text and standalone."""
        from backend.services.slide_generator import generate_slide_content

        captured_prompt = []

        def capture_generate_content(model=None, contents=None, config=None, **kwargs):
            # contents is a list of {"role": ..., "parts": [{"text": ...}]}
            prompt_text = contents[0]["parts"][0]["text"] if contents else ""
            captured_prompt.append(prompt_text)
            return MagicMock(
                text=json.dumps(SAMPLE_SLIDE_CONTENT),
                candidates=[MagicMock(finish_reason=MagicMock(name="STOP"))],
                usage_metadata=MagicMock(prompt_token_count=10, candidates_token_count=5),
            )

        with patch('backend.services.llm_adapter.gemini_adapter.genai') as mock_genai:
            mock_client = MagicMock()
            mock_client.models.generate_content.side_effect = capture_generate_content
            mock_genai.Client.return_value = mock_client

            generate_slide_content(
                content="Test content", subject="Math", grade="7",
                title="Test", api_key="fake", deck_format="detailed",
            )

        assert captured_prompt, "generate_content was not called"
        assert "DETAILED DECK" in captured_prompt[0]

    def test_presenter_format_includes_talking_points(self):
        """Presenter format prompt should mention key talking points."""
        from backend.services.slide_generator import generate_slide_content

        captured_prompt = []

        def capture_generate_content(model=None, contents=None, config=None, **kwargs):
            prompt_text = contents[0]["parts"][0]["text"] if contents else ""
            captured_prompt.append(prompt_text)
            return MagicMock(
                text=json.dumps(SAMPLE_SLIDE_CONTENT),
                candidates=[MagicMock(finish_reason=MagicMock(name="STOP"))],
                usage_metadata=MagicMock(prompt_token_count=10, candidates_token_count=5),
            )

        with patch('backend.services.llm_adapter.gemini_adapter.genai') as mock_genai:
            mock_client = MagicMock()
            mock_client.models.generate_content.side_effect = capture_generate_content
            mock_genai.Client.return_value = mock_client

            generate_slide_content(
                content="Test content", subject="Math", grade="7",
                title="Test", api_key="fake", deck_format="presenter",
            )

        assert captured_prompt, "generate_content was not called"
        assert "PRESENTER SLIDES" in captured_prompt[0]
